"""
Update bubble chart data in place on slide 16, preserving frame size /
position / styling done in PowerPoint. Replaces series with top 8 bundles
ranked by growth %.

Usage:
  .venv/bin/python update_chart_data.py
"""
import csv
from pathlib import Path

from pptx import Presentation
from pptx.chart.data import BubbleChartData
from pptx.dml.color import RGBColor
from lxml import etree

HERE = Path(__file__).parent
APR_CSV = Path("/Users/gblanc/git_clones/genai/prompt-registry/lib/analytics-output/hub-analytics-2026-04-23T13-50-37-by-bundle.csv")
MAY_CSV = HERE / "hub-analytics-2026-05-21T15-42-23-by-bundle.csv"
PPTX_FILE = HERE / "LEAP_Iteration_Review_26.2.3_updated.pptx"

TOP_N = 8

SERIES_DEFS = [
    (">1000% growth", RGBColor(0xE6, 0x39, 0x46), lambda r: r["growth"] > 1000),
    ("200–1000%",     RGBColor(0xF4, 0xA2, 0x61), lambda r: 200 < r["growth"] <= 1000),
    ("50–200%",       RGBColor(0x2A, 0x9D, 0x8F), lambda r: 50 < r["growth"] <= 200),
    ("<50%",          RGBColor(0x7B, 0x9A, 0xD0), lambda r: r["growth"] <= 50),
]

# Manual label-position overrides to break overlap in lower-left cluster.
# Valid values for bubble dLblPos: ctr, l, r, t, b. Default = r.
LABEL_POSITION = {
    "foursight-code-assessment":     "t",
    "opentelemetry-tracing-toolkit": "l",
    "refx-nre-qa-agents":            "b",
    "amadeus-microservice-coding-guidebook": "l",
}


def load(csv_path):
    out = {}
    with open(csv_path) as f:
        for row in csv.DictReader(f):
            out[row["Bundle ID"]] = int(row["Total Downloads"])
    return out


def build_rows():
    apr = load(APR_CSV)
    may = load(MAY_CSV)
    rows = []
    for bid, may_dl in may.items():
        if bid not in apr:
            continue
        apr_dl = apr[bid]
        if apr_dl <= 0:
            continue
        delta = may_dl - apr_dl
        if delta <= 0:
            continue
        rows.append({
            "id": bid,
            "may": may_dl,
            "apr": apr_dl,
            "delta": delta,
            "growth": (delta / apr_dl) * 100,
        })
    rows.sort(key=lambda r: r["growth"], reverse=True)
    return rows[:TOP_N]


NS_C = "{http://schemas.openxmlformats.org/drawingml/2006/chart}"
NS_A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"


def reapply_series_styling(chart, series_rows):
    """After replace_data, reset fill colors + per-point bundle-name labels."""
    plot = chart.plots[0]
    for s_idx, series in enumerate(plot.series):
        label, color, members = series_rows[s_idx]
        # Fill
        fill = series.format.fill
        fill.solid()
        fill.fore_color.rgb = color
        solidFill = fill._xPr.find(f"{NS_A}solidFill")
        if solidFill is not None:
            srgb = solidFill.find(f"{NS_A}srgbClr")
            if srgb is not None:
                # transparency 30%
                if srgb.find(f"{NS_A}alpha") is None:
                    alpha = etree.SubElement(srgb, f"{NS_A}alpha")
                    alpha.set("val", "70000")
        # Outline
        line = series.format.line
        line.color.rgb = RGBColor(0x33, 0x33, 0x33)
        from pptx.util import Pt
        line.width = Pt(0.5)

        # Per-point labels
        ser_xml = series._element
        for dl in ser_xml.findall(f"{NS_C}dLbls"):
            ser_xml.remove(dl)
        dLbls = etree.SubElement(ser_xml, f"{NS_C}dLbls")
        for i, r in enumerate(members):
            dLbl = etree.SubElement(dLbls, f"{NS_C}dLbl")
            idx = etree.SubElement(dLbl, f"{NS_C}idx")
            idx.set("val", str(i))
            # Schema order: idx, layout?, tx?, numFmt?, spPr?, txPr?,
            # dLblPos?, showLegendKey, showVal, ...
            tx = etree.SubElement(dLbl, f"{NS_C}tx")
            rich = etree.SubElement(tx, f"{NS_C}rich")
            bodyPr = etree.SubElement(rich, f"{NS_A}bodyPr")
            bodyPr.set("wrap", "square")
            etree.SubElement(rich, f"{NS_A}lstStyle")
            p = etree.SubElement(rich, f"{NS_A}p")
            run = etree.SubElement(p, f"{NS_A}r")
            rPr = etree.SubElement(run, f"{NS_A}rPr")
            rPr.set("sz", "700")
            t = etree.SubElement(run, f"{NS_A}t")
            t.text = r["id"]
            pos = LABEL_POSITION.get(r["id"], "r")
            dLblPos = etree.SubElement(dLbl, f"{NS_C}dLblPos")
            dLblPos.set("val", pos)
            for tag in ("showLegendKey", "showVal", "showCatName",
                        "showSerName", "showPercent", "showBubbleSize"):
                el = etree.SubElement(dLbl, f"{NS_C}{tag}")
                el.set("val", "0")
        for tag in ("showLegendKey", "showVal", "showCatName",
                    "showSerName", "showPercent", "showBubbleSize"):
            el = etree.SubElement(dLbls, f"{NS_C}{tag}")
            el.set("val", "0")


def main():
    rows = build_rows()
    print(f"top {TOP_N} by growth%:")
    for r in rows:
        print(f"  {r['id']}: {r['growth']:.1f}% (Δ{r['delta']}, {r['may']} dl)")

    prs = Presentation(str(PPTX_FILE))
    slide = prs.slides[15]
    chart = None
    for sh in slide.shapes:
        if sh.has_chart and sh.chart.chart_type.name == "BUBBLE":
            chart = sh.chart
            break
    if chart is None:
        raise SystemExit("bubble chart not found on slide 16")

    chart_data = BubbleChartData()
    series_rows = []
    for label, color, predicate in SERIES_DEFS:
        members = [r for r in rows if predicate(r)]
        series_rows.append((label, color, members))
        s = chart_data.add_series(label)
        for r in members:
            s.add_data_point(r["may"], r["growth"], r["delta"])

    chart.replace_data(chart_data)
    reapply_series_styling(chart, series_rows)

    prs.save(str(PPTX_FILE))
    print(f"updated {PPTX_FILE}")
    print("series counts:", [(l, len(m)) for l, _, m in series_rows])


if __name__ == "__main__":
    main()
