"""
Build native PowerPoint bubble chart on slide 16, replacing the stale
"5 FASTEST GROWING BUNDLES" table. Chart uses embedded data so user can
edit values directly in PowerPoint after the file is opened.

Axes:
  X = total downloads as of May 21 (current size)
  Y = % growth Apr 23 -> May 21 (log scale)
  Bubble size = absolute download delta
  4 color-coded series by growth bucket (>1000%, 200-1000%, 50-200%, <50%)
"""
import csv
from copy import deepcopy
from pathlib import Path

from pptx import Presentation
from pptx.chart.data import BubbleChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.util import Emu, Inches, Pt
from lxml import etree

HERE = Path(__file__).parent
APR_CSV = Path("/Users/gblanc/git_clones/genai/prompt-registry/lib/analytics-output/hub-analytics-2026-04-23T13-50-37-by-bundle.csv")
MAY_CSV = HERE / "hub-analytics-2026-05-21T15-42-23-by-bundle.csv"
PPTX_IN = HERE / "LEAP_Iteration_Review_26.2.3.pptx"
PPTX_OUT = HERE / "LEAP_Iteration_Review_26.2.3_updated.pptx"


def load(csv_path):
    out = {}
    with open(csv_path) as f:
        for row in csv.DictReader(f):
            out[row["Bundle ID"]] = int(row["Total Downloads"])
    return out


TOP_N = 8


def build_data():
    apr = load(APR_CSV)
    may = load(MAY_CSV)
    rows = []
    for bid, may_dl in may.items():
        if bid not in apr:
            continue
        apr_dl = apr[bid]
        if apr_dl <= 0:
            continue
        delta = may_dl - apr_dl
        if delta <= 0:
            continue
        growth = (delta / apr_dl) * 100
        rows.append({
            "id": bid,
            "may": may_dl,
            "apr": apr_dl,
            "delta": delta,
            "growth": growth,
        })
    rows.sort(key=lambda r: r["growth"], reverse=True)
    return rows[:TOP_N]


SERIES_DEFS = [
    (">1000% growth", RGBColor(0xE6, 0x39, 0x46), lambda r: r["growth"] > 1000),
    ("200–1000%",     RGBColor(0xF4, 0xA2, 0x61), lambda r: 200 < r["growth"] <= 1000),
    ("50–200%",       RGBColor(0x2A, 0x9D, 0x8F), lambda r: 50 < r["growth"] <= 200),
    ("<50%",          RGBColor(0x7B, 0x9A, 0xD0), lambda r: r["growth"] <= 50),
]


def remove_shapes_by_name(slide, names):
    """Delete shapes whose .name matches names (set)."""
    spTree = slide.shapes._spTree
    for shape in list(slide.shapes):
        if shape.name in names:
            spTree.remove(shape._element)


# Names of shapes that compose the "5 FASTEST GROWING BUNDLES" table block.
# Derived from inspect_slide16.py output (indices 24..49).
TABLE_SHAPE_NAMES = {
    "Text 34", "Text 35",
    "Shape 36", "Text 37", "Text 38", "Text 39",
    "Shape 40", "Text 41", "Text 42", "Text 43",
    "Shape 44", "Text 45", "Text 46", "Text 47",
    "Shape 48", "Text 49", "Text 50", "Text 51",
    "Shape 52", "Text 53", "Text 54", "Text 55",
    "Shape 56", "Text 57", "Text 58", "Text 59",
    # Stray duplicate '820' label sitting inside chart bbox (left of original
    # table). Shape 'Text 21' at L=7.40,T=3.19 — collides with new chart.
    "Text 21",
}

# Bundles to label on chart. Labelling all 23 produces unreadable overlap.
# Keep outliers and the largest movers; the rest stay legible via legend +
# tooltip in PowerPoint edit mode.
LABEL_WHITELIST = {
    "amadeus-microservice-coding-guidebook",
    "airline-solutions-architecture",
    "archreview",
    "workflow-nevio",
    "skubedocs",
    "pr-review",
    "java-engineering",
    "cnadocs",
    "log-verbosity-reduction",
    "task-driven-agents",
    "reflex",
    "foursight-code-assessment",
}


def add_data_labels(plot, rows_for_series):
    """Attach bundle-name data labels to each point in a bubble series."""
    ser = plot.series[0]  # one series per plot here? actually we add many
    # Not used; we set labels per-series in main flow.


def main():
    rows = build_data()
    rows.sort(key=lambda r: r["growth"], reverse=True)
    print(f"{len(rows)} bundles plotted")

    prs = Presentation(str(PPTX_IN))
    slide = prs.slides[15]

    # Remove stale table shapes
    remove_shapes_by_name(slide, TABLE_SHAPE_NAMES)

    # Build bubble chart data: 4 series for color-coding & legend
    chart_data = BubbleChartData()
    series_rows = []
    for label, color, predicate in SERIES_DEFS:
        members = [r for r in rows if predicate(r)]
        series_rows.append((label, color, members))
        s = chart_data.add_series(label)
        for r in members:
            # x = current downloads, y = growth %, size = delta
            s.add_data_point(r["may"], r["growth"], r["delta"])

    # Insert chart on right side of slide
    x = Inches(6.47)
    y = Inches(2.66)
    cx = Inches(6.50)   # extend a bit closer to slide right edge
    cy = Inches(2.85)

    graphic_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.BUBBLE, x, y, cx, cy, chart_data
    )
    chart = graphic_frame.chart
    chart.has_title = True
    chart.chart_title.text_frame.text = "Bundle adoption momentum (Apr 23 → May 21)"
    for run in chart.chart_title.text_frame.paragraphs[0].runs:
        run.font.size = Pt(11)
        run.font.bold = True

    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(8)

    # Color each series + add bundle-name labels
    for plot in chart.plots:
        plot.has_data_labels = True
        for s_idx, series in enumerate(plot.series):
            label, color, members = series_rows[s_idx]
            # Series fill color
            fill = series.format.fill
            fill.solid()
            fill.fore_color.rgb = color
            # Slight transparency via XML (python-pptx has no helper)
            solidFill = fill._xPr.find(
                "{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill"
            )
            if solidFill is not None:
                srgb = solidFill.find(
                    "{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr"
                )
                if srgb is not None:
                    alpha = etree.SubElement(
                        srgb,
                        "{http://schemas.openxmlformats.org/drawingml/2006/main}alpha",
                    )
                    alpha.set("val", "70000")  # 70 %
            # Outline: thin dark
            line = series.format.line
            line.color.rgb = RGBColor(0x33, 0x33, 0x33)
            line.width = Pt(0.5)

            # Per-point data labels = bundle ID
            ns = "{http://schemas.openxmlformats.org/drawingml/2006/chart}"
            ser_xml = series._element  # c:ser
            # Remove any existing dLbls so we control fully
            for dl in ser_xml.findall(f"{ns}dLbls"):
                ser_xml.remove(dl)
            dLbls = etree.SubElement(ser_xml, f"{ns}dLbls")
            for i, r in enumerate(members):
                if r["id"] not in LABEL_WHITELIST:
                    continue
                dLbl = etree.SubElement(dLbls, f"{ns}dLbl")
                idx = etree.SubElement(dLbl, f"{ns}idx")
                idx.set("val", str(i))
                tx = etree.SubElement(dLbl, f"{ns}tx")
                rich = etree.SubElement(
                    tx,
                    "{http://schemas.openxmlformats.org/drawingml/2006/chart}rich",
                )
                bodyPr = etree.SubElement(
                    rich,
                    "{http://schemas.openxmlformats.org/drawingml/2006/main}bodyPr",
                )
                bodyPr.set("wrap", "square")
                etree.SubElement(
                    rich,
                    "{http://schemas.openxmlformats.org/drawingml/2006/main}lstStyle",
                )
                p = etree.SubElement(
                    rich,
                    "{http://schemas.openxmlformats.org/drawingml/2006/main}p",
                )
                run = etree.SubElement(
                    p,
                    "{http://schemas.openxmlformats.org/drawingml/2006/main}r",
                )
                rPr = etree.SubElement(
                    run,
                    "{http://schemas.openxmlformats.org/drawingml/2006/main}rPr",
                )
                rPr.set("sz", "700")
                t = etree.SubElement(
                    run,
                    "{http://schemas.openxmlformats.org/drawingml/2006/main}t",
                )
                t.text = r["id"]
                showLegendKey = etree.SubElement(dLbl, f"{ns}showLegendKey")
                showLegendKey.set("val", "0")
                showVal = etree.SubElement(dLbl, f"{ns}showVal")
                showVal.set("val", "0")
                showCatName = etree.SubElement(dLbl, f"{ns}showCatName")
                showCatName.set("val", "0")
                showSerName = etree.SubElement(dLbl, f"{ns}showSerName")
                showSerName.set("val", "0")
                showPercent = etree.SubElement(dLbl, f"{ns}showPercent")
                showPercent.set("val", "0")
                showBubbleSize = etree.SubElement(dLbl, f"{ns}showBubbleSize")
                showBubbleSize.set("val", "0")
            # Series-level: hide all default fields
            for tag in (
                "showLegendKey", "showVal", "showCatName",
                "showSerName", "showPercent", "showBubbleSize",
            ):
                el = etree.SubElement(dLbls, f"{ns}{tag}")
                el.set("val", "0")

    # Axes: log Y, linear X, with titles
    cat = chart.category_axis  # x in bubble = "category"? actually value
    val = chart.value_axis     # y
    # python-pptx maps bubble's x-axis as category_axis; both are valueAx
    val.has_major_gridlines = True
    cat.has_major_gridlines = False

    # Titles
    cat.has_title = True
    cat.axis_title.text_frame.text = "Total downloads (May 21)"
    for run in cat.axis_title.text_frame.paragraphs[0].runs:
        run.font.size = Pt(8)
    val.has_title = True
    val.axis_title.text_frame.text = "Growth % (log)"
    for run in val.axis_title.text_frame.paragraphs[0].runs:
        run.font.size = Pt(8)

    # Set Y axis to logarithmic via XML, with explicit min=1 max=10000 so
    # >1000% outliers are visible (default auto-scale capped at 1000).
    nsc = "{http://schemas.openxmlformats.org/drawingml/2006/chart}"

    def force_log_scale(axis_el, lo, hi):
        scaling = axis_el.find(f"{nsc}scaling")
        if scaling is None:
            scaling = etree.SubElement(axis_el, f"{nsc}scaling")
        # Wipe and rebuild scaling children in the order required by schema:
        # logBase, orientation, max, min
        for child in list(scaling):
            scaling.remove(child)
        logBase = etree.SubElement(scaling, f"{nsc}logBase")
        logBase.set("val", "10")
        orient = etree.SubElement(scaling, f"{nsc}orientation")
        orient.set("val", "minMax")
        mx = etree.SubElement(scaling, f"{nsc}max")
        mx.set("val", str(hi))
        mn = etree.SubElement(scaling, f"{nsc}min")
        mn.set("val", str(lo))

    force_log_scale(val._element, 1, 10000)

    # Tick label font size
    for axis in (cat, val):
        try:
            axis.tick_labels.font.size = Pt(8)
        except Exception:
            pass

    prs.save(str(PPTX_OUT))
    print(f"saved {PPTX_OUT}")
    print("series counts:", [(l, len(m)) for l, _, m in series_rows])


if __name__ == "__main__":
    main()
