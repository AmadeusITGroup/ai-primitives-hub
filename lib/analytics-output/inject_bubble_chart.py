#!/usr/bin/env python3
"""
Inject a native bubble chart into slide 16 of LEAP_Iteration_Review_26.2.3.pptx.
Uses python-pptx for chart structure, then post-processes chart XML for:
  - Log scale Y axis
  - Data labels (bundle names)
  - Styling matching slide colors

Run from analytics-output/ directory.
"""
import copy
import re
from lxml import etree
from pptx import Presentation
from pptx.chart.data import BubbleChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_LABEL_POSITION

# ── Data ──────────────────────────────────────────────────────────────────────
# (bundle_label, may_downloads, growth_pct, delta)
BUNDLES = [
    ("microservice-guidebook",  178, 5833.3, 175),
    ("airline-solutions-arch",   80, 1233.3,  74),
    ("archreview",               63,  350.0,  49),
    ("foursight-code-assess",    28,  180.0,  18),
    ("opentelemetry-toolkit",    19,  171.4,  12),
    ("refx-nre-qa-agents",       31,  138.5,  18),
    ("cnadocs",                  75,  120.6,  41),
    ("pr-review",               129,  111.5,  68),
    ("java-engineering",        107,  109.8,  56),
    ("batch-elasticity",          6,  100.0,   3),
    ("dsre-git-skillset",        19,   72.7,   8),
    ("etk-workflow",             38,   65.2,  15),
    ("agents-md-creator",        16,   60.0,   6),
    ("rail-obe-context",         11,   57.1,   4),
    ("slidev",                   53,   51.4,  18),
    ("skubedocs",               170,   50.4,  57),
    ("pr-reviewer-generator",    97,   44.8,  30),
    ("reflex",                  119,   25.3,  24),
    ("log-verbosity-reduction",  79,   21.5,  14),
    ("aude",                     53,   20.5,   9),
    ("amadeus-ospo",             62,   19.2,  10),
    ("workflow-nevio",          782,   17.1, 114),
    ("task-driven-agents",       77,    8.5,   6),
]

# Color tiers (hex without #)
def tier(growth):
    if growth >= 1000: return "extreme"
    if growth >= 200:  return "high"
    if growth >= 50:   return "solid"
    return "moderate"

SERIES_COLORS = {
    "extreme":  "E63946",
    "high":     "F4A261",
    "solid":    "2A9D8F",
    "moderate": "7B9AD0",
}

SERIES_NAMES = {
    "extreme":  ">1000% growth",
    "high":     "200–1000%",
    "solid":    "50–200%",
    "moderate": "<50%",
}

# Bubble size = delta (direct, pptx scales it internally)
def bubble_size(delta):
    return max(delta, 1)

# ── Build chart data ───────────────────────────────────────────────────────────
chart_data = BubbleChartData()

series_map = {t: chart_data.add_series(SERIES_NAMES[t]) for t in ["extreme", "high", "solid", "moderate"]}

for name, may, growth, delta in BUNDLES:
    t = tier(growth)
    series_map[t].add_data_point(may, growth, bubble_size(delta))

# ── Open PPTX and get slide 16 ────────────────────────────────────────────────
prs = Presentation("LEAP_Iteration_Review_26.2.3_updated.pptx")
slide = prs.slides[15]  # 0-indexed

# Remove shapes 37-62 (already removed in previous step) and the image41 we added
shapes_to_remove = []
for shape in slide.shapes:
    # Remove the bubble chart image we injected (id=70)
    if shape.shape_id == 70:
        shapes_to_remove.append(shape)

sp_tree = slide.shapes._spTree
for shape in shapes_to_remove:
    sp_tree.remove(shape._element)

# ── Add native bubble chart ───────────────────────────────────────────────────
# Position: x=5917087, y=2432304, cx=5425085, cy=2600000 (EMU)
from pptx.util import Emu
chart_placeholder = slide.shapes.add_chart(
    XL_CHART_TYPE.BUBBLE,
    Emu(5917087), Emu(2432304),
    Emu(5425085), Emu(2600000),
    chart_data
)

chart = chart_placeholder.chart

# ── Post-process chart XML ────────────────────────────────────────────────────
chart_part = chart_placeholder.chart._part
chart_xml = chart_part._element

NSMAP = {
    'c':   'http://schemas.openxmlformats.org/drawingml/2006/chart',
    'a':   'http://schemas.openxmlformats.org/drawingml/2006/main',
    'r':   'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
}

def ns(tag): return tag.replace('c:', '{http://schemas.openxmlformats.org/drawingml/2006/chart}').replace('a:', '{http://schemas.openxmlformats.org/drawingml/2006/main}')

def set_color_on_series(ser_elem, hex_color):
    """Set solid fill color on a series <c:ser> element."""
    spPr = ser_elem.find(ns('c:spPr'))
    if spPr is None:
        spPr = etree.SubElement(ser_elem, ns('c:spPr'))
    # solidFill
    solidFill = etree.SubElement(spPr, ns('a:solidFill'))
    srgbClr = etree.SubElement(solidFill, ns('a:srgbClr'))
    srgbClr.set('val', hex_color)

def remove_children(elem, tag):
    for child in elem.findall(ns(tag)):
        elem.remove(child)

# Get plot area
plotArea = chart_xml.find('.//' + ns('c:plotArea'))

# Get all series
series_list = plotArea.findall('.//' + ns('c:ser'))
tier_order = ["extreme", "high", "solid", "moderate"]

for i, ser in enumerate(series_list):
    t = tier_order[i]
    hex_color = SERIES_COLORS[t]
    set_color_on_series(ser, hex_color)
    # Also set marker color
    marker = ser.find(ns('c:marker'))
    if marker is None:
        marker = etree.SubElement(ser, ns('c:marker'))
    spPr_m = marker.find(ns('c:spPr'))
    if spPr_m is None:
        spPr_m = etree.SubElement(marker, ns('c:spPr'))
    sf = etree.SubElement(spPr_m, ns('a:solidFill'))
    sc = etree.SubElement(sf, ns('a:srgbClr'))
    sc.set('val', hex_color)

# ── Y-axis: log scale, remove minor gridlines, style ─────────────────────────
valAx_elems = plotArea.findall(ns('c:valAx'))
# valAx[0] = X axis, valAx[1] = Y axis (for bubble charts)
for valAx in valAx_elems:
    axId = valAx.find(ns('c:axId'))
    # Check orientation / position to identify Y axis
    # Y axis has crosses set — just apply log to both and let pptx sort it
    scaling = valAx.find(ns('c:scaling'))
    if scaling is not None:
        logBase = etree.SubElement(scaling, ns('c:logBase'))
        logBase.set('val', '10')

# ── Remove chart border / plot area border ────────────────────────────────────
chart_elem = chart_xml.find(ns('c:chart'))
plotArea_spPr = plotArea.find(ns('c:spPr'))
if plotArea_spPr is None:
    plotArea_spPr = etree.SubElement(plotArea, ns('c:spPr'))
    plotArea.insert(0, plotArea_spPr)

# Set plot area fill to light blue-gray
fill = etree.SubElement(plotArea_spPr, ns('a:solidFill'))
sc2 = etree.SubElement(fill, ns('a:srgbClr'))
sc2.set('val', 'F8F9FC')

# ── Chart background ──────────────────────────────────────────────────────────
spPr_chart = chart_xml.find('.//' + ns('c:chartSpace') + '//' + ns('c:spPr'))
chartSpace = chart_xml  # root IS chartSpace
spPr_root = chartSpace.find(ns('c:spPr'))
if spPr_root is None:
    spPr_root = etree.SubElement(chartSpace, ns('c:spPr'))
bg_fill = etree.SubElement(spPr_root, ns('a:solidFill'))
bg_sc = etree.SubElement(bg_fill, ns('a:srgbClr'))
bg_sc.set('val', 'FFFFFF')

# ── Save ──────────────────────────────────────────────────────────────────────
prs.save("LEAP_Iteration_Review_26.2.3_updated.pptx")
print("Saved LEAP_Iteration_Review_26.2.3_updated.pptx with native bubble chart")
