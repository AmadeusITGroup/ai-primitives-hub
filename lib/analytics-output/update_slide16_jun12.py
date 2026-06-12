#!/usr/bin/env python3
"""
Update slide 16 of LEAP_Iteration_Review_26.2.3.pptx with Jun 12 analytics.
Prev snapshot: hub-analytics-2026-05-21T15-42-23
Curr snapshot: hub-analytics-2026-06-12T09-38-29

Run:
  python3.14 update_slide16_jun12.py
"""
import csv
import copy
from pathlib import Path
from lxml import etree
from pptx import Presentation
from pptx.chart.data import BubbleChartData
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt, Emu

HERE = Path(__file__).parent

PREV_BUNDLE_CSV = HERE / "hub-analytics-2026-05-21T15-42-23-by-bundle.csv"
CURR_BUNDLE_CSV = HERE / "hub-analytics-2026-06-12T09-38-29-by-bundle.csv"
PREV_SOURCE_CSV = HERE / "hub-analytics-2026-05-21T15-42-23-by-source.csv"
CURR_SOURCE_CSV = HERE / "hub-analytics-2026-06-12T09-38-29-by-source.csv"

SRC_PPTX  = HERE / "LEAP_Iteration_Review_26.2.3.pptx"
OUT_PPTX  = HERE / "LEAP_Iteration_Review_26.2.3_updated.pptx"

PREV_DATE = "May 21"
CURR_DATE = "Jun 12"

TOP_N = 8

SERIES_DEFS = [
    (">1000% growth", RGBColor(0xE6, 0x39, 0x46), lambda r: r["growth"] > 1000),
    ("200–1000%",     RGBColor(0xF4, 0xA2, 0x61), lambda r: 200 < r["growth"] <= 1000),
    ("50–200%",       RGBColor(0x2A, 0x9D, 0x8F), lambda r: 50 < r["growth"] <= 200),
    ("<50%",          RGBColor(0x7B, 0x9A, 0xD0), lambda r: r["growth"] <= 50),
]

LABEL_POSITION = {
    "foursight-code-assessment":           "t",
    "opentelemetry-tracing-toolkit":       "l",
    "refx-nre-qa-agents":                  "b",
    "amadeus-microservice-coding-guidebook": "l",
    "nevio-solution-architecture-runway":  "t",
    "specifications-as-code-collection":   "b",
}

NS_C = "{http://schemas.openxmlformats.org/drawingml/2006/chart}"
NS_A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"


# ── Data helpers ───────────────────────────────────────────────────────────────

def load_bundles(csv_path):
    out = {}
    with open(csv_path) as f:
        for row in csv.DictReader(f):
            out[row["Bundle ID"]] = int(row["Total Downloads"])
    return out


def load_source_ids(csv_path):
    ids = set()
    with open(csv_path) as f:
        for row in csv.DictReader(f):
            ids.add(row["Source ID"])
    return ids


def build_bubble_rows(prev, curr):
    rows = []
    for bid, curr_dl in curr.items():
        if bid not in prev:
            continue
        prev_dl = prev[bid]
        if prev_dl <= 0:
            continue
        delta = curr_dl - prev_dl
        if delta <= 0:
            continue
        rows.append({
            "id":     bid,
            "curr":   curr_dl,
            "prev":   prev_dl,
            "delta":  delta,
            "growth": (delta / prev_dl) * 100,
        })
    rows.sort(key=lambda r: r["growth"], reverse=True)
    return rows[:TOP_N]


def pct(curr, prev):
    d = round((curr - prev) / prev * 100)
    arrow = "↑" if d >= 0 else "↓"
    sign  = "+" if d >= 0 else ""
    return f"{sign}{d}% {arrow}"


# ── Slide text helpers ─────────────────────────────────────────────────────────

def set_text_run(shape, new_text):
    """Replace first non-empty run's text; preserve all formatting."""
    txBody = shape.text_frame._txBody
    for p_el in txBody.findall(f"{NS_A}p"):
        for r in p_el.findall(f"{NS_A}r"):
            t = r.find(f"{NS_A}t")
            if t is not None:
                t.text = new_text
                return


def replace_run_containing(shape, needle, replacement):
    """Find the run whose text contains needle and patch it in place."""
    txBody = shape.text_frame._txBody
    for p_el in txBody.findall(f"{NS_A}p"):
        for r in p_el.findall(f"{NS_A}r"):
            t = r.find(f"{NS_A}t")
            if t is not None and needle in (t.text or ""):
                t.text = t.text.replace(needle, replacement)
                return


def find_shape_by_name(slide, name, exclude_id=None):
    """Find shape by name, searching groups too. Optionally exclude a shape_id."""
    for sh in slide.shapes:
        if sh.name == name and (exclude_id is None or sh.shape_id != exclude_id):
            return sh
        if sh.shape_type == 6:  # GROUP
            for child in sh.shapes:
                if child.name == name and (exclude_id is None or child.shape_id != exclude_id):
                    return child
    return None


def remove_shape_by_id(slide, shape_id):
    sp_tree = slide.shapes._spTree
    for sh in slide.shapes:
        if sh.shape_id == shape_id:
            sp_tree.remove(sh._element)
            return True
    return False


# ── Bubble chart styling ───────────────────────────────────────────────────────

def reapply_series_styling(chart, series_rows):
    plot = chart.plots[0]
    for s_idx, series in enumerate(plot.series):
        label, color, members = series_rows[s_idx]
        fill = series.format.fill
        fill.solid()
        fill.fore_color.rgb = color
        solidFill = fill._xPr.find(f"{NS_A}solidFill")
        if solidFill is not None:
            srgb = solidFill.find(f"{NS_A}srgbClr")
            if srgb is not None and srgb.find(f"{NS_A}alpha") is None:
                alpha = etree.SubElement(srgb, f"{NS_A}alpha")
                alpha.set("val", "70000")
        line = series.format.line
        line.color.rgb = RGBColor(0x33, 0x33, 0x33)
        line.width = Pt(0.5)

        ser_xml = series._element
        for dl in ser_xml.findall(f"{NS_C}dLbls"):
            ser_xml.remove(dl)
        dLbls = etree.SubElement(ser_xml, f"{NS_C}dLbls")
        for i, r in enumerate(members):
            dLbl   = etree.SubElement(dLbls, f"{NS_C}dLbl")
            idx    = etree.SubElement(dLbl,  f"{NS_C}idx")
            idx.set("val", str(i))
            tx     = etree.SubElement(dLbl,  f"{NS_C}tx")
            rich   = etree.SubElement(tx,    f"{NS_C}rich")
            bodyPr = etree.SubElement(rich,  f"{NS_A}bodyPr")
            bodyPr.set("wrap", "square")
            etree.SubElement(rich, f"{NS_A}lstStyle")
            p_el   = etree.SubElement(rich,  f"{NS_A}p")
            run    = etree.SubElement(p_el,  f"{NS_A}r")
            rPr    = etree.SubElement(run,   f"{NS_A}rPr")
            rPr.set("sz", "700")
            t_el   = etree.SubElement(run,   f"{NS_A}t")
            t_el.text = r["id"]
            pos = LABEL_POSITION.get(r["id"], "r")
            dLblPos = etree.SubElement(dLbl, f"{NS_C}dLblPos")
            dLblPos.set("val", pos)
            for tag in ("showLegendKey", "showVal", "showCatName",
                        "showSerName", "showPercent", "showBubbleSize"):
                el = etree.SubElement(dLbl, f"{NS_C}{tag}")
                el.set("val", "0")
        for tag in ("showLegendKey", "showVal", "showCatName",
                    "showSerName", "showPercent", "showBubbleSize"):
            el = etree.SubElement(dLbls, f"{NS_C}{tag}")
            el.set("val", "0")


# ── Main ───────────────────────────────────────────────────────────────────────

def main():
    prev_bundles = load_bundles(PREV_BUNDLE_CSV)
    curr_bundles = load_bundles(CURR_BUNDLE_CSV)
    prev_sources = load_source_ids(PREV_SOURCE_CSV)
    curr_sources = load_source_ids(CURR_SOURCE_CSV)

    # KPI values
    prev_dl  = sum(prev_bundles.values())
    curr_dl  = sum(curr_bundles.values())
    prev_src = len(prev_sources)
    curr_src = len(curr_sources)
    prev_bun = len(prev_bundles)
    curr_bun = len(curr_bundles)

    # New teams = sources in curr not in prev
    new_teams   = len(curr_sources - prev_sources)
    # New bundles = bundle IDs in curr not in prev
    new_bun_ids = set(curr_bundles) - set(prev_bundles)
    new_bun_cnt = len(new_bun_ids)

    print(f"Downloads: {prev_dl} → {curr_dl}  {pct(curr_dl, prev_dl)}")
    print(f"Sources:   {prev_src} → {curr_src}  {pct(curr_src, prev_src)}")
    print(f"Bundles:   {prev_bun} → {curr_bun}  {pct(curr_bun, prev_bun)}")
    print(f"New teams: {new_teams}  |  New bundles: {new_bun_cnt}")

    # Top 5 bundles
    top5 = sorted(curr_bundles.items(), key=lambda x: x[1], reverse=True)[:5]
    print("Top 5 bundles:", top5)

    # Bubble chart rows
    bubble_rows = build_bubble_rows(prev_bundles, curr_bundles)
    print(f"\nBubble chart top {TOP_N} by growth:")
    for r in bubble_rows:
        print(f"  {r['id']}: {r['growth']:.1f}% (Δ{r['delta']}, curr={r['curr']})")

    # ── Open presentation ──────────────────────────────────────────────────────
    prs   = Presentation(str(SRC_PPTX))
    slide = prs.slides[15]

    # ── Title ──────────────────────────────────────────────────────────────────
    title_sh = find_shape_by_name(slide, "Text Placeholder 5")
    if title_sh:
        replace_run_containing(title_sh, f"from {PREV_DATE} to May 21",
                               f"from {PREV_DATE} to {CURR_DATE}")
        # also handle case where May 21 is not there yet
        replace_run_containing(title_sh, "evolutions from Apr 23 to May 21",
                               f"evolutions from {PREV_DATE} to {CURR_DATE}")
    else:
        print("WARNING: title shape not found")

    # ── Header KPIs ───────────────────────────────────────────────────────────
    dl_val  = find_shape_by_name(slide, "Text 5")
    dl_pct  = find_shape_by_name(slide, "Text 6")
    src_val = find_shape_by_name(slide, "Text 8")
    src_pct = find_shape_by_name(slide, "Text 9")
    bun_val = find_shape_by_name(slide, "Text 11")
    bun_pct = find_shape_by_name(slide, "Text 12")

    if dl_val:  set_text_run(dl_val,  f"{curr_dl:,}")
    if dl_pct:  replace_run_containing(dl_pct, "+59%", pct(curr_dl, prev_dl).split()[0])
    if src_val: set_text_run(src_val, str(curr_src))
    if src_pct: replace_run_containing(src_pct, "+16%", pct(curr_src, prev_src).split()[0])
    if bun_val: set_text_run(bun_val, str(curr_bun))
    if bun_pct: replace_run_containing(bun_pct, "+20%", pct(curr_bun, prev_bun).split()[0])

    print(f"\nKPIs updated: dl={curr_dl:,} {pct(curr_dl,prev_dl)}, "
          f"src={curr_src} {pct(curr_src,prev_src)}, "
          f"bun={curr_bun} {pct(curr_bun,prev_bun)}")

    # ── Remove stray Text 24 (id=64, overlaps chart area) ─────────────────────
    removed = remove_shape_by_id(slide, 64)
    if removed:
        print("Removed stray Text 24 (id=64) '820'")

    # ── Top 5 bar list ────────────────────────────────────────────────────────
    # Name shapes: Text 19, Text 22, Text 25, Text 28, Text 31
    # Count shapes: Text 24 (id=27), Text 27, Text 30, Text 33  (no Text 21)
    # Bar shapes: Shape 20, Shape 23, Shape 26, Shape 29, Shape 32

    name_shape_names  = ["Text 19", "Text 22", "Text 25", "Text 28", "Text 31"]
    count_shape_names = ["Text 24", "Text 27", "Text 30", "Text 33"]
    bar_shape_names   = ["Shape 20", "Shape 23", "Shape 26", "Shape 29", "Shape 32"]

    top1_dl   = top5[0][1]
    base_width = 2331720  # EMU for the top bar (full width)

    for i, (bid, dl) in enumerate(top5):
        # Name
        name_sh = find_shape_by_name(slide, name_shape_names[i])
        if name_sh:
            set_text_run(name_sh, bid)

        # Count (Text 21 doesn't exist; Text 24 = count for bundle 2 → index 1)
        # counts are for bundles 2–5 (indices 1–4)
        if i >= 1:
            cnt_sh = find_shape_by_name(slide, count_shape_names[i - 1],
                                        exclude_id=64)  # skip stray
            if cnt_sh:
                set_text_run(cnt_sh, str(dl))

        # Bar width
        bar_sh = find_shape_by_name(slide, bar_shape_names[i])
        if bar_sh:
            new_w = int(base_width * dl / top1_dl)
            bar_sh.width = new_w

    print("Top 5 bars updated.")

    # ── Bubble chart ──────────────────────────────────────────────────────────
    chart = None
    for sh in slide.shapes:
        if sh.has_chart and sh.chart.chart_type.name == "BUBBLE":
            chart = sh.chart
            break
    if chart is None:
        print("WARNING: bubble chart not found on slide 16 — skipping chart update")
    else:
        chart_data = BubbleChartData()
        series_rows = []
        for label, color, predicate in SERIES_DEFS:
            members = [r for r in bubble_rows if predicate(r)]
            series_rows.append((label, color, members))
            s = chart_data.add_series(label)
            for r in members:
                s.add_data_point(r["curr"], r["growth"], r["delta"])

        chart.replace_data(chart_data)
        reapply_series_styling(chart, series_rows)
        print(f"Bubble chart updated. Series counts: "
              f"{[(l, len(m)) for l, _, m in series_rows]}")

    # ── Save ──────────────────────────────────────────────────────────────────
    prs.save(str(OUT_PPTX))
    print(f"\nSaved: {OUT_PPTX}")


if __name__ == "__main__":
    main()
