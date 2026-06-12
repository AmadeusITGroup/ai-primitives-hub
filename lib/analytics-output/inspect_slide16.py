from pptx import Presentation
from pptx.util import Emu

p = Presentation("LEAP_Iteration_Review_26.2.3.pptx")
print(f"slides: {len(p.slides)}")
slide = p.slides[15]  # 16th slide
print(f"slide16 title placeholder check")
for i, shape in enumerate(slide.shapes):
    print(f"  [{i}] type={shape.shape_type} name='{shape.name}' "
          f"L={Emu(shape.left).inches:.2f} T={Emu(shape.top).inches:.2f} "
          f"W={Emu(shape.width).inches:.2f} H={Emu(shape.height).inches:.2f} "
          f"has_table={shape.has_table} has_chart={shape.has_chart}")
    if shape.has_text_frame:
        txt = shape.text_frame.text[:100].replace("\n", " | ")
        print(f"      text: {txt}")
    if shape.has_table:
        t = shape.table
        print(f"      table: {len(t.rows)}x{len(t.columns)}")
        for r in t.rows:
            for c in r.cells:
                print(f"        cell: {c.text[:60]}")
